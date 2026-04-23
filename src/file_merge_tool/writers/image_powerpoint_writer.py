from __future__ import annotations

from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
GOLD = RGBColor(143, 111, 62)
CHARCOAL = RGBColor(38, 34, 29)
MUTED = RGBColor(96, 87, 75)
IVORY = RGBColor(251, 248, 241)


def write_image_powerpoint(
    path: Path,
    *,
    title: str,
    header: dict[str, Any],
    summary: dict[str, Any],
    items: list[dict[str, Any]],
    skipped_items: list[dict[str, Any]],
    warnings: list[dict[str, Any]],
) -> Path:
    path.parent.mkdir(parents=True, exist_ok=True)
    presentation = Presentation()
    presentation.slide_width = SLIDE_WIDTH
    presentation.slide_height = SLIDE_HEIGHT

    _add_header_slide(
        presentation,
        title=title,
        header=header,
        summary=summary,
        skipped_items=skipped_items,
        warnings=warnings,
    )
    for index, item in enumerate(items, start=1):
        _add_image_slide(presentation, index=index, item=item)

    presentation.save(path)
    return path


def _add_header_slide(
    presentation: Presentation,
    *,
    title: str,
    header: dict[str, Any],
    summary: dict[str, Any],
    skipped_items: list[dict[str, Any]],
    warnings: list[dict[str, Any]],
) -> None:
    slide = _blank_slide(presentation)
    _set_background(slide)
    _add_top_rule(slide)

    _text_box(
        slide,
        title,
        left=Inches(0.7),
        top=Inches(0.58),
        width=Inches(11.9),
        height=Inches(0.48),
        font_size=24,
        bold=True,
        color=CHARCOAL,
    )
    _text_box(
        slide,
        "Merge header",
        left=Inches(0.72),
        top=Inches(1.1),
        width=Inches(3.2),
        height=Inches(0.24),
        font_size=8,
        color=GOLD,
        uppercase=True,
    )

    lines = [
        ("Schema", header.get("schema") or header.get("schema_name") or ""),
        ("Job ID", header.get("job_id") or ""),
        ("Kind", header.get("kind") or ""),
        ("Classification", header.get("classification") or ""),
        ("Generated", header.get("generated_at") or ""),
        ("Source root", header.get("source_root") or ""),
        ("Setting", header.get("setting_name") or ""),
        ("Traversal", (header.get("traversal") or {}).get("order", "")),
        ("Follow symlinks", str((header.get("traversal") or {}).get("follow_symlinks", False))),
        ("Exclude folders", ", ".join((header.get("exclude") or {}).get("folders", []))),
        ("Exclude extensions", ", ".join((header.get("exclude") or {}).get("extensions", []))),
        ("Exclude files", ", ".join((header.get("exclude") or {}).get("file_names", []))),
        (
            "Sensitivity markers",
            ", ".join(
                [
                    *((header.get("sensitivity") or {}).get("default_markers", [])),
                    *((header.get("sensitivity") or {}).get("additional_markers", [])),
                ]
            ),
        ),
        ("Items", str(summary.get("item_count", 0))),
        ("Skipped", str(summary.get("skipped_count", 0))),
        ("Read errors", str(summary.get("error_skipped_count", 0))),
        ("Warnings", str(summary.get("warning_count", 0))),
    ]
    _key_value_list(
        slide,
        lines,
        left=Inches(0.72),
        top=Inches(1.5),
        width=Inches(12),
        height=Inches(4.6),
    )

    issue_lines = _issue_lines(skipped_items, warnings)
    if issue_lines:
        _text_box(
            slide,
            "Skipped / warnings",
            left=Inches(0.72),
            top=Inches(6.18),
            width=Inches(2.6),
            height=Inches(0.25),
            font_size=8,
            color=GOLD,
            uppercase=True,
        )
        _text_box(
            slide,
            "\n".join(issue_lines[:5]),
            left=Inches(0.72),
            top=Inches(6.45),
            width=Inches(11.8),
            height=Inches(0.7),
            font_size=8,
            color=MUTED,
        )


def _add_image_slide(presentation: Presentation, *, index: int, item: dict[str, Any]) -> None:
    slide = _blank_slide(presentation)
    _set_background(slide)
    _add_top_rule(slide)

    title = f"{index:03d}  {item.get('relative_path', '')}"
    _text_box(
        slide,
        title,
        left=Inches(0.58),
        top=Inches(0.36),
        width=Inches(12.2),
        height=Inches(0.38),
        font_size=14,
        bold=True,
        color=CHARCOAL,
    )

    image_path = Path(str(item["absolute_path"]))
    _add_fitted_picture(
        slide,
        image_path=image_path,
        source_width=int(item.get("width") or 1),
        source_height=int(item.get("height") or 1),
        box_left=Inches(0.65),
        box_top=Inches(1.0),
        box_width=Inches(9.0),
        box_height=Inches(5.82),
    )

    meta_lines = [
        ("Absolute path", str(item.get("absolute_path") or "")),
        ("Modified", str(item.get("modified_at") or "")),
        ("MIME", str(item.get("mime_type") or "")),
        ("Pixels", f"{item.get('width') or ''} x {item.get('height') or ''}"),
        ("Size", f"{item.get('file_size_bytes') or 0:,} bytes"),
    ]
    sensitivity = item.get("sensitivity") or {}
    matched = ", ".join(sensitivity.get("matched_markers") or [])
    if matched:
        meta_lines.append(("Matched markers", matched))

    _text_box(
        slide,
        "Metadata",
        left=Inches(9.95),
        top=Inches(1.05),
        width=Inches(2.4),
        height=Inches(0.25),
        font_size=8,
        color=GOLD,
        uppercase=True,
    )
    _key_value_list(
        slide,
        meta_lines,
        left=Inches(9.95),
        top=Inches(1.4),
        width=Inches(2.75),
        height=Inches(5.1),
        font_size=8,
    )


def _blank_slide(presentation: Presentation):
    return presentation.slides.add_slide(presentation.slide_layouts[6])


def _set_background(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = IVORY


def _add_top_rule(slide) -> None:
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.55),
        Inches(0.18),
        Inches(12.2),
        Inches(0.035),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = GOLD
    shape.line.fill.background()


def _add_fitted_picture(
    slide,
    *,
    image_path: Path,
    source_width: int,
    source_height: int,
    box_left,
    box_top,
    box_width,
    box_height,
) -> None:
    ratio = min(box_width / source_width, box_height / source_height)
    width = int(source_width * ratio)
    height = int(source_height * ratio)
    left = int(box_left + (box_width - width) / 2)
    top = int(box_top + (box_height - height) / 2)
    slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)


def _key_value_list(
    slide,
    items: list[tuple[str, str]],
    *,
    left,
    top,
    width,
    height,
    font_size: int = 8,
) -> None:
    text = "\n".join(f"{label}: {value}" for label, value in items)
    _text_box(
        slide,
        text,
        left=left,
        top=top,
        width=width,
        height=height,
        font_size=font_size,
        color=CHARCOAL,
    )


def _text_box(
    slide,
    text: str,
    *,
    left,
    top,
    width,
    height,
    font_size: int,
    color,
    bold: bool = False,
    uppercase: bool = False,
) -> None:
    shape = slide.shapes.add_textbox(left, top, width, height)
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.LEFT
    run = paragraph.add_run()
    run.text = text.upper() if uppercase else text
    run.font.name = "Yu Gothic"
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color


def _issue_lines(
    skipped_items: list[dict[str, Any]],
    warnings: list[dict[str, Any]],
) -> list[str]:
    lines: list[str] = []
    for item in skipped_items[:3]:
        lines.append(f"Skipped: {item.get('relative_path', '')} ({item.get('reason', '')})")
    for item in warnings[:3]:
        lines.append(f"Warning: {item.get('relative_path', '')} ({item.get('reason', '')})")
    return lines
