from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation

from file_merge_tool.application.merge_image import merge_image
from file_merge_tool.domain.config import ExcludeConfig, MergeRequest


def test_image_merge_writes_normal_and_sensitive_html(tmp_path: Path) -> None:
    root = tmp_path / "root"
    root.mkdir()
    _write_png(root / "normal.png")
    _write_png(root / f"secret_\u6a5f\u5bc6.png")

    request = MergeRequest(
        root_path=root,
        output_dir=tmp_path / "out",
        output_name="images.html",
        output_stem="images",
        exclude=ExcludeConfig(),
        kind="image-merge",
        image_output_formats=("html",),
    )

    result = merge_image(request)

    assert result.item_count == 2
    assert len(result.output_paths) == 2
    normal_html = result.output_paths[0].read_text(encoding="utf-8")
    sensitive_html = result.output_paths[1].read_text(encoding="utf-8")
    assert result.output_paths[0].name == "images_\u30de\u30fc\u30b8.html"
    assert result.output_paths[1].name == "\u6a5f\u5bc6_images_\u30de\u30fc\u30b8.html"
    assert "normal.png" in normal_html
    assert "secret_" not in normal_html
    assert "secret_" in sensitive_html
    assert "data:image/png;base64," in normal_html
    assert '"schema": "file-merge-tool/image-html-report/v1"' in normal_html


def test_image_merge_records_unreadable_image_warning(tmp_path: Path) -> None:
    root = tmp_path / "root"
    root.mkdir()
    (root / "broken.png").write_text("not an image", encoding="utf-8")

    request = MergeRequest(
        root_path=root,
        output_dir=tmp_path / "out",
        output_name="images.html",
        output_stem="images",
        exclude=ExcludeConfig(),
        kind="image-merge",
        image_output_formats=("html",),
    )

    result = merge_image(request)

    assert result.item_count == 0
    assert result.error_skipped_count == 1
    output = result.output_paths[0].read_text(encoding="utf-8")
    assert "broken.png" in output
    assert "read_error" in output


def test_image_merge_writes_powerpoint_outputs(tmp_path: Path) -> None:
    root = tmp_path / "root"
    root.mkdir()
    _write_png(root / "normal.png")
    _write_png(root / f"secret_\u6a5f\u5bc6.png")

    request = MergeRequest(
        root_path=root,
        output_dir=tmp_path / "out",
        output_name="images.pptx",
        output_stem="images",
        exclude=ExcludeConfig(),
        kind="image-merge",
        image_output_formats=("pptx",),
    )

    result = merge_image(request)

    assert result.item_count == 2
    assert [path.name for path in result.output_paths] == [
        "images_\u30de\u30fc\u30b8.pptx",
        "\u6a5f\u5bc6_images_\u30de\u30fc\u30b8.pptx",
    ]
    normal_deck = Presentation(result.output_paths[0])
    sensitive_deck = Presentation(result.output_paths[1])
    assert len(normal_deck.slides) == 2
    assert len(sensitive_deck.slides) == 2


def _write_png(path: Path) -> None:
    image = Image.new("RGB", (3, 2), color=(120, 80, 40))
    image.save(path)
